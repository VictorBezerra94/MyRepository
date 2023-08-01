import matplotlib.pyplot as plt
import matplotlib_venn
import numpy as np
import math
import xlrd
import pandas as pd
import statistics as st
import random
import pptx
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from scipy import stats

def cos(a = -3,b = 3, A=1):
    xs = list ([*(x/10.0 for x in range (a*10,b*10+1))])
    plt.plot(xs, [A*math.cos(x) for x in xs], '-')
    plt.show()

def parabola(a = -3, b = 3):
    xs = list([*(x/10 for x in range (10*a,10*b+1))])
    plt.plot( xs, [x**2 for x in xs], '--' , label = 'x^2')
    plt.plot( xs, [x**3 for x in xs], '.' ,  label = 'x^3')
    plt.plot( xs, [x**4 for x in xs], '.-' ,  label = 'x^4')
    plt.legend()
    plt.show()

#tempo em meses
def juros(i, t):
    xs = list(range(0,t+1))
    base = 1+i/100
    plt.plot(xs, [math.pow(base,t) for t in xs], '--', label = 'Juros compostos')
    plt.plot(xs, [1+i/100*t        for t in xs], '^-',  label = 'Juros simples')
    plt.legend()
    plt.show()

def circle(a=0, b=0 , R=1):
    xv = list([*(x/100 for x in range(100*a-R*100, 100*(R)+100*a+1   ))])
    plt.plot(xv, [ b + math.sqrt( (-(x-a)**2 + R**2) ) for x in xv]    , color= 'b')
    plt.plot(xv, [ b - math.sqrt( (-(x-a)**2 + R**2) ) for x in xv]    , color= 'b')
    plt.plot(xv, [b for x in xv]                                       , color= 'g', label='eixo x')
    plt.plot([a for x in range(-R,R+1) ] , [y for y in range (-R+b,R+1+b)] , color= 'g', label='eixo y')
    plt.legend()
    plt.show()

def ven_diagram():

    v = matplotlib_venn.venn2       (subsets=(3,4,7))
    matplotlib_venn.venn2_circles   (subsets=(3,4,7) , linestyle="dashed", linewidth=2)

    #Costomizacao do circulo A
    v.get_patch_by_id('10').set_alpha(1.0)
    v.get_patch_by_id('10').set_facecolor('pink')
    v.get_patch_by_id('10').set_edgecolor('black')
    v.get_patch_by_id('10').set_linestyle('-') # {'-', '--', '-.', ':', '', (offset, on-off-seq), ...}
    v.get_patch_by_id('10').set_linewidth(5)
    v.get_patch_by_id('10').set_hatch('/') # {'/', '\', '|', '-', '+', 'x', 'o', 'O', '.', '*'}
    #v.get_label_by_id('10').set_text('A')
    v.get_label_by_id('A').set_text('A')

    #Costomizacao do circulo B
    v.get_patch_by_id('01').set_alpha(0.1)
    v.get_patch_by_id('01').set_facecolor('darkblue')
    v.get_patch_by_id('01').set_edgecolor('black')
    v.get_patch_by_id('01').set_linestyle('--')
    v.get_patch_by_id('01').set_linewidth(5)
    v.get_patch_by_id('01').set_hatch('/') #
   #v.get_label_by_id('01').set_text('A')
    v.get_label_by_id('B').set_text('B')

    #Costomizacao do circulo C
    v.get_patch_by_id('11').set_alpha(0.1)
    v.get_patch_by_id('11').set_facecolor('seagreen')
   #v.get_patch_by_id('11').set_edgecolor('M
    v.get_patch_by_id('11').set_linestyle('--')
    v.get_patch_by_id('11').set_linewidth(5)
    v.get_patch_by_id('11').set_hatch('/') #

    plt.show()
    plt.savefig('Figure1.png')

def make_PP():

    presentation = pptx.Presentation()
    slide1 = presentation.slides.add_slide(presentation.slide_layouts[6])
    x       = Inches(2)
    y       = Inches(1.5)
    altura  = Inches(4)
    largura = Inches(6)

    pic = slide1.shapes.add_picture('Figure_1.png',x,y,largura,altura)

    x       = Inches(1)
    y       = Inches(0.5)
    altura  = Inches(1)
    largura = Inches(8)

    caixa_de_texto = slide1.shapes.add_textbox(x,y,largura,altura)

    text_frame = caixa_de_texto.text_frame
    paragraph = text_frame.add_paragraph()
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.text = 'Venn Diagram'
    paragraph.font.size = Pt(38)
    paragraph.font.bold = True
    presentation.save("meu_PPT2.pptx")


def acum_aux(A):
    size = len(A)
    acum = 0
    list=[]
    for i in range(0,size):
        acum = acum + A[i]
        list.append(acum)
    return list

def acum(x):
    xs = list([*(a for a in range(0, 31))])
    plt.plot(xs, acum_aux(xs))
    plt.show()

def dobro():
    xs = list([*(a for a in range(0, 31))])
    dob = plt.plot(xs, [2**x/100 for x in xs])
    plt.grid(axis='y')
    plt.show()

acum(2)